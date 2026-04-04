"""
Document Processor — text extraction for non-PDF file formats.
Supports: .docx, .pptx, .txt, .md, .epub
"""

import logging
import os
import re
from html.parser import HTMLParser
from typing import Any, Callable, Dict, List

logger = logging.getLogger(__name__)

# Type alias — every extractor returns a list of these
PageData = Dict[str, Any]
# keys: text, page_number, file_name, file_path, total_pages

SUPPORTED_EXTENSIONS = frozenset({".docx", ".pptx", ".txt", ".md", ".epub"})


# ---------------------------------------------------------------------------
# Shared helper
# ---------------------------------------------------------------------------

def _make_page(
    text: str,
    page_number: int,
    file_path: str,
    total_pages: int,
) -> PageData:
    """Single source of truth for the page-data schema."""
    return {
        "text": text,
        "page_number": page_number,
        "file_name": os.path.basename(file_path),
        "file_path": file_path,
        "total_pages": total_pages,
    }


# ---------------------------------------------------------------------------
# Public entry point
# ---------------------------------------------------------------------------

def extract_text_from_file(file_path: str) -> List[PageData]:
    """
    Extract text from a supported non-PDF file.

    Returns:
        List of page dicts.  Empty list for unsupported / empty / corrupt files.
    """
    ext = os.path.splitext(file_path)[1].lower()

    handlers: Dict[str, Callable[[str], List[PageData]]] = {
        ".docx": _from_docx,
        ".pptx": _from_pptx,
        ".txt":  _from_text,
        ".md":   _from_text,
        ".epub": _from_epub,
    }

    handler = handlers.get(ext)
    if handler is None:
        logger.warning("Unsupported file format '%s': %s", ext, file_path)
        return []

    try:
        return handler(file_path)
    except Exception:
        logger.exception("Failed to extract text from %s", file_path)
        return []


# ---------------------------------------------------------------------------
# Format-specific extractors
# ---------------------------------------------------------------------------

def _from_docx(file_path: str) -> List[PageData]:
    """Extract text from a .docx (treated as a single logical page)."""
    from docx import Document

    doc = Document(file_path)
    text = "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())
    if not text.strip():
        return []
    return [_make_page(text, page_number=1, file_path=file_path, total_pages=1)]


def _from_pptx(file_path: str) -> List[PageData]:
    """Extract text from a .pptx (one page per slide)."""
    from pptx import Presentation

    prs = Presentation(file_path)
    total_slides = len(prs.slides)
    pages: List[PageData] = []

    for i, slide in enumerate(prs.slides, start=1):
        text = "\n".join(
            shape.text
            for shape in slide.shapes
            if hasattr(shape, "text") and shape.text.strip()
        )
        if text.strip():
            pages.append(
                _make_page(text, page_number=i, file_path=file_path, total_pages=total_slides)
            )
    return pages


def _from_text(file_path: str) -> List[PageData]:
    """Extract text from a plain-text or markdown file."""
    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        text = f.read()
    if not text.strip():
        return []
    return [_make_page(text, page_number=1, file_path=file_path, total_pages=1)]


# ---- EPUB helpers --------------------------------------------------------

class _HTMLTextExtractor(HTMLParser):
    """Minimal HTML→text converter that inserts whitespace around block tags."""

    _BLOCK_TAGS = frozenset({
        "p", "div", "h1", "h2", "h3", "h4", "h5", "h6",
        "li", "blockquote", "br", "tr", "section", "article",
    })

    def __init__(self) -> None:
        super().__init__()
        self.parts: List[str] = []

    def handle_starttag(self, tag: str, attrs: Any) -> None:
        if tag.lower() in self._BLOCK_TAGS:
            self.parts.append("\n")

    def handle_data(self, data: str) -> None:
        self.parts.append(data)

    def get_text(self) -> str:
        raw = "".join(self.parts)
        return re.sub(r"\s+", " ", raw).strip()


def _from_epub(file_path: str) -> List[PageData]:
    """Extract text from an .epub (one page per document item)."""
    import ebooklib
    from ebooklib import epub

    book = epub.read_epub(file_path)
    items = list(book.get_items_of_type(ebooklib.ITEM_DOCUMENT))
    total_items = len(items)

    pages: List[PageData] = []
    for i, item in enumerate(items, start=1):
        parser = _HTMLTextExtractor()
        parser.feed(item.get_content().decode("utf-8", errors="replace"))
        text = parser.get_text()
        if text:
            pages.append(
                _make_page(text, page_number=i, file_path=file_path, total_pages=total_items)
            )
    return pages